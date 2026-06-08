"""
Build one or more Band It Layer deck slides as .pptx (1950s newspaper style).

Usage:
  python docs/presentations/band-it-layer-deck/scripts/build_slide.py 1
  python docs/presentations/band-it-layer-deck/scripts/build_slide.py 1 2 3
  python docs/presentations/band-it-layer-deck/scripts/build_slide.py --all
"""

from __future__ import annotations

import argparse
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SLIDES_DIR = ROOT / "slides"
ASSETS_DIR = ROOT / "assets"
OUTPUT_DIR = ROOT

# Band It newspaper tokens
PAPER = RGBColor(0xFA, 0xFA, 0xFA)
INK = RGBColor(0x0A, 0x0A, 0x0A)
MUTED = RGBColor(0x5C, 0x5C, 0x5C)
ACCENT = RGBColor(0x00, 0x14, 0xFF)
NEON = RGBColor(0xFF, 0xDD, 0x00)
RULE_SOFT = RGBColor(0xE0, 0xE0, 0xE0)
PANEL_FILL = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEADLINE = "Arial Black"
FONT_BODY = "Georgia"
FONT_KICKER = "Consolas"


def parse_slide_md(path: Path) -> dict:
    text = path.read_text(encoding="utf-8")
    headline_m = re.search(r"^## Headline\s*\n(.+?)(?=\n## )", text, re.M | re.S)
    kicker_m = re.search(r"^## Kicker.*\n(.+?)(?=\n## )", text, re.M | re.S)
    body_m = re.search(r"^## Body\s*\n(.+?)(?=\n## )", text, re.M | re.S)
    layout_m = re.search(r"^## Layout\s*\n(.+?)(?=\n## )", text, re.M | re.S)

    headline = (headline_m.group(1).strip() if headline_m else path.stem)
    kicker = (kicker_m.group(1).strip() if kicker_m else "")
    body_raw = body_m.group(1).strip() if body_m else ""
    layout = (layout_m.group(1).strip().lower() if layout_m else "")
    image_only = "image-only" in layout or "image only" in layout
    no_image = "no-image" in layout or "no image" in layout
    asset_m = re.search(r"^## Asset\s*\n(.+?)(?=\n## |\Z)", text, re.M | re.S)

    paragraphs: list[str] = []
    bullets: list[str] = []
    for block in re.split(r"\n\n+", body_raw):
        block = block.strip()
        if not block:
            continue
        lines = [ln.strip() for ln in block.split("\n") if ln.strip()]
        if all(ln.startswith("- ") for ln in lines):
            bullets.extend(ln[2:] for ln in lines)
        else:
            paragraphs.append(" ".join(lines))

    slug = path.stem
    num = slug.split("-", 1)[0]
    asset = None
    if asset_m:
        raw = asset_m.group(1).strip().strip("`")
        if raw and not raw.startswith("("):
            candidate = (path.parent / raw).resolve()
            if candidate.exists():
                asset = candidate
    if asset is None:
        n = int(num)
        for name in (f"slide {n}.png", f"slide {n}.jpg", f"slide{n}.png"):
            candidate = ASSETS_DIR / name
            if candidate.exists():
                asset = candidate
                break
    if asset is None and not no_image and not image_only:
        matches = sorted(
            list(ASSETS_DIR.glob(f"slide-{num}-*.jpg"))
            + list(ASSETS_DIR.glob(f"slide-{num}-*.jpeg"))
            + list(ASSETS_DIR.glob(f"slide-{num}-*.png"))
        )
        asset = matches[-1] if matches else None
    if asset is None and image_only:
        n = int(num)
        for name in (f"slide {n}.png", f"slide {n}.jpg"):
            candidate = ASSETS_DIR / name
            if candidate.exists():
                asset = candidate
                break

    caption_m = re.search(r"^## (?:Image caption|Fig caption)\s*\n(.+?)(?=\n## )", text, re.M | re.S)
    fig_caption = caption_m.group(1).strip() if caption_m else f"Fig. {int(num)} — Signal desk"

    return {
        "num": num,
        "headline": headline,
        "kicker": kicker,
        "paragraphs": paragraphs,
        "bullets": bullets,
        "asset": asset,
        "fig_caption": fig_caption,
        "no_image": no_image,
        "image_only": image_only,
    }


def _add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _style_run(run, *, size: int, bold: bool = False, color=INK, font: str = FONT_BODY):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _filled_rect(slide, left, top, width, height, fill, *, line_color=None, line_pt=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line_color is not None:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_pt)
    else:
        sh.line.fill.background()
    return sh


def _hline(slide, left, top, width, *, thick=Pt(1), color=INK):
    return _filled_rect(slide, left, top, width, thick, color)


def _vline(slide, left, top, height, *, thick=Pt(1), color=RULE_SOFT):
    return _filled_rect(slide, left, top, thick, height, color)


def _add_masthead(slide, margin, content_w, data: dict) -> None:
    """Double-rule masthead band + newspaper page folio."""
    top = Inches(0.38)
    _hline(slide, margin, top, content_w, thick=Pt(2), color=INK)
    _hline(slide, margin, top + Pt(5), content_w, thick=Pt(1), color=INK)
    _hline(slide, margin, top + Pt(9), content_w, thick=Pt(1), color=RULE_SOFT)

    page_n = int(data["num"])
    folio_w = Inches(1.1)
    folio = _add_textbox(slide, margin + content_w - folio_w, Inches(0.5), folio_w, Inches(0.32))
    folio.text_frame.clear()
    fp = folio.text_frame.paragraphs[0]
    fr = fp.add_run()
    fr.text = f"Page {page_n}"
    _style_run(fr, size=10, bold=False, color=INK, font=FONT_BODY)
    fp.alignment = PP_ALIGN.RIGHT


def _add_signal_desk_panel(
    slide,
    left,
    top,
    width,
    height,
    *,
    fig_caption: str,
    slide_num: str,
) -> None:
    """Editorial 'signal vs noise' schematic — lines and borders only."""
    # Outer double frame
    _filled_rect(slide, left, top, width, height, PANEL_FILL, line_color=INK, line_pt=1)
    inset = Inches(0.06)
    inner_l = left + inset
    inner_t = top + inset
    inner_w = width - inset * 2
    inner_h = height - inset * 2
    _filled_rect(
        slide, inner_l, inner_t, inner_w, inner_h, PAPER, line_color=RULE_SOFT, line_pt=1
    )

    pad = Inches(0.14)
    ix = inner_l + pad
    iy = inner_t + pad
    iw = inner_w - pad * 2

    # Panel kicker
    lbl = _add_textbox(slide, ix, iy, iw, Inches(0.28))
    lbl.text_frame.clear()
    lp = lbl.text_frame.paragraphs[0]
    lr = lp.add_run()
    lr.text = "SIGNAL DESK · INCOMING"
    _style_run(lr, size=7, color=MUTED, font=FONT_KICKER)
    _hline(slide, ix, iy + Inches(0.3), iw, thick=Pt(1), color=INK)

    # Noise: staggered ticker lines (muted)
    noise_top = iy + Inches(0.42)
    noise_h = inner_h * 0.52
    line_gap = Pt(11)
    widths = (1.0, 0.92, 0.78, 0.88, 0.65, 0.95, 0.72, 0.85, 0.6, 0.9, 0.7, 0.82, 0.55, 0.88)
    y = noise_top
    for i, frac in enumerate(widths):
        lw = int(iw * frac)
        color = RULE_SOFT if i % 3 else MUTED
        _hline(slide, ix, y, lw, thick=Pt(1), color=color)
        y += line_gap

    # Divider + label
    mid_y = noise_top + noise_h
    _hline(slide, ix, mid_y, iw, thick=Pt(2), color=INK)

    tag = _add_textbox(slide, ix, mid_y + Inches(0.08), iw, Inches(0.22))
    tag.text_frame.clear()
    tp = tag.text_frame.paragraphs[0]
    tr = tp.add_run()
    tr.text = "▸ ONE SIGNAL"
    _style_run(tr, size=8, bold=True, color=ACCENT, font=FONT_KICKER)

    # Highlighted signal bar (neon rule + blue fill strip)
    sig_y = mid_y + Inches(0.34)
    sig_h = Inches(0.42)
    _filled_rect(slide, ix, sig_y, iw, sig_h, RGBColor(0xF5, 0xF6, 0xFF), line_color=ACCENT, line_pt=2)
    _filled_rect(slide, ix, sig_y, Pt(5), sig_h, NEON)

    sig_txt = _add_textbox(slide, ix + Inches(0.12), sig_y + Pt(4), iw - Inches(0.14), sig_h - Pt(6))
    sig_txt.text_frame.clear()
    sp = sig_txt.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = "Budget approved · expansion signal"
    _style_run(sr, size=9, bold=True, color=INK, font=FONT_BODY)

    # Brass dials (circles) + crosshatch corner lines
    dial_y = sig_y + Inches(0.65)
    dial_r = Inches(0.2)
    for dx in (0, iw * 0.38, iw * 0.76):
        dial = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, int(ix + dx), int(dial_y), int(dial_r), int(dial_r)
        )
        dial.fill.solid()
        dial.fill.fore_color.rgb = PAPER
        dial.line.color.rgb = INK
        dial.line.width = Pt(1)
        # tick mark
        cx = ix + dx + dial_r / 2
        _hline(slide, cx - Pt(8), dial_y + dial_r / 2 - Pt(1), Pt(16), thick=Pt(1), color=INK)

    # Static corner hatch (diagonal rules)
    hx = ix + iw - Inches(0.55)
    hy = dial_y + Inches(0.05)
    for i in range(6):
        off = i * Pt(6)
        sh = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            int(hx + off),
            int(hy + off),
            Pt(1),
            Inches(0.35),
        )
        sh.fill.solid()
        sh.fill.fore_color.rgb = RULE_SOFT
        sh.line.fill.background()
        sh.rotation = -35

    # Fig caption
    cap_y = inner_t + inner_h - Inches(0.38)
    cap = _add_textbox(slide, ix, cap_y, iw, Inches(0.32))
    cap.text_frame.clear()
    cp = cap.text_frame.paragraphs[0]
    cr = cp.add_run()
    cr.text = fig_caption
    _style_run(cr, size=7, color=MUTED, font=FONT_KICKER)
    cp.alignment = PP_ALIGN.CENTER


def _add_image_panel(slide, data: dict, img_left, img_w, sh) -> None:
    frame_top = Inches(1.0)
    frame_h = sh - frame_top - Inches(0.55)
    pad = Inches(0.08)
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left, frame_top, img_w, frame_h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = PANEL_FILL
    frame.line.color.rgb = INK
    frame.line.width = Pt(1)

    slide.shapes.add_picture(
        str(data["asset"]),
        img_left + pad,
        frame_top + pad,
        width=img_w - pad * 2,
        height=frame_h - pad * 2,
    )


def add_image_only_slide(prs: Presentation, data: dict) -> None:
    """Full-bleed hero slide — generated art only, no copy."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER

    sw = prs.slide_width
    sh = prs.slide_height
    if not data.get("asset") or not data["asset"].exists():
        raise FileNotFoundError(f"Image-only slide {data['num']} requires an asset image")

    slide.shapes.add_picture(str(data["asset"]), 0, 0, width=sw, height=sh)


def add_slide(prs: Presentation, data: dict) -> None:
    if data.get("image_only"):
        add_image_only_slide(prs, data)
        return

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER

    sw = prs.slide_width
    sh = prs.slide_height
    margin = Inches(0.55)
    content_w = sw - margin * 2

    _add_masthead(slide, margin, content_w, data)

    # Kicker
    if data["kicker"]:
        kbox = _add_textbox(slide, margin, Inches(0.9), content_w - Inches(0.7), Inches(0.35))
        ktf = kbox.text_frame
        ktf.clear()
        p = ktf.paragraphs[0]
        r = p.add_run()
        r.text = data["kicker"].upper()
        _style_run(r, size=9, color=MUTED, font=FONT_KICKER)
        p.alignment = PP_ALIGN.LEFT

    text_left = margin
    text_w = int(content_w * 0.56)
    panel_left = text_left + text_w + Inches(0.22)
    panel_w = content_w - text_w - Inches(0.22)

    # Column rule
    _vline(slide, panel_left - Inches(0.1), Inches(0.95), sh - Inches(1.35), thick=Pt(1), color=RULE_SOFT)

    # Yellow accent + headline block
    accent = _filled_rect(slide, text_left, Inches(1.05), Pt(4), Inches(1.55), NEON)
    hbox = _add_textbox(
        slide, text_left + Inches(0.14), Inches(1.0), text_w - Inches(0.14), Inches(1.65)
    )
    htf = hbox.text_frame
    htf.word_wrap = True
    htf.clear()
    hp = htf.paragraphs[0]
    hr = hp.add_run()
    hr.text = data["headline"]
    _style_run(hr, size=22, bold=True, font=FONT_HEADLINE)
    hp.line_spacing = 1.05
    _hline(slide, text_left, Inches(2.62), text_w, thick=Pt(2), color=INK)
    _hline(slide, text_left, Inches(2.68), text_w, thick=Pt(1), color=RULE_SOFT)

    # Body inset frame
    body_top = Inches(2.82)
    body_h = sh - body_top - Inches(0.62)
    _filled_rect(
        slide,
        text_left,
        body_top,
        text_w,
        body_h,
        PANEL_FILL,
        line_color=RULE_SOFT,
        line_pt=1,
    )
    _filled_rect(slide, text_left, body_top, Pt(3), body_h, INK)

    bbox = _add_textbox(
        slide, text_left + Inches(0.12), body_top + Inches(0.1), text_w - Inches(0.2), body_h - Inches(0.15)
    )
    btf = bbox.text_frame
    btf.word_wrap = True
    btf.vertical_anchor = MSO_ANCHOR.TOP
    btf.clear()
    first = True

    for para in data["paragraphs"]:
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = para
        _style_run(r, size=13, font=FONT_BODY)
        p.space_after = Pt(10)
        p.line_spacing = 1.15

    for bullet in data["bullets"]:
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        p.level = 0
        r = p.add_run()
        r.text = f"—  {bullet}"
        _style_run(r, size=12, font=FONT_BODY)
        p.space_after = Pt(4)
        p.line_spacing = 1.1

    # Right column: image or signal-desk schematic
    use_image = (
        not data.get("no_image")
        and data.get("asset")
        and data["asset"].exists()
    )
    if use_image:
        _add_image_panel(slide, data, panel_left, panel_w, sh)
    else:
        _add_signal_desk_panel(
            slide,
            panel_left,
            Inches(1.0),
            panel_w,
            Inches(5.85),
            fig_caption=data.get("fig_caption", f"Fig. {data['num']} — Signal desk"),
            slide_num=data["num"],
        )

def build(slide_nums: list[str], out_name: str) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for num in slide_nums:
        matches = list(SLIDES_DIR.glob(f"{num.zfill(2)}-*.md"))
        if not matches:
            matches = list(SLIDES_DIR.glob(f"{num}-*.md"))
        if not matches:
            raise FileNotFoundError(f"No slide markdown for number {num}")
        add_slide(prs, parse_slide_md(matches[0]))

    out = OUTPUT_DIR / out_name
    prs.save(out)
    return out


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("slides", nargs="*", help="Slide numbers e.g. 1 2 3")
    parser.add_argument("--all", action="store_true", help="Build all slides in slides/")
    args = parser.parse_args()

    if args.all:
        nums = sorted({p.name.split("-", 1)[0] for p in SLIDES_DIR.glob("*.md") if p.name != "README.md"})
        out_name = "band-it-layer-deck.pptx"
    elif args.slides:
        nums = args.slides
        if len(nums) == 1:
            out_name = "band-it-layer-deck-slide-{:02d}.pptx".format(int(nums[0]))
        elif len(nums) >= 2:
            hi = max(int(n) for n in nums)
            lo = min(int(n) for n in nums)
            out_name = f"band-it-layer-deck-slides-{lo:02d}-{hi:02d}.pptx"
        else:
            out_name = "band-it-layer-deck.pptx"
    else:
        parser.error("Provide slide numbers or --all")

    path = build(nums, out_name)
    print(path)


if __name__ == "__main__":
    main()
